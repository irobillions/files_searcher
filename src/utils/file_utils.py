import re
import os
import docx
import win32com.client
import openpyxl
import xlrd
from pptx import Presentation
from PyPDF2 import PdfReader
import pytesseract
from pdf2image import convert_from_path
import camelot
import comtypes.client
import io

def get_filename()-> str:
    return input("Donnez la cible de recherche: ").strip()


def get_boolean_input(prompt):
    while True:
        response = input(prompt).strip().lower()
        if response in {"oui", "non"}:
            return response == "oui"
        else:
            print("Réponse non valide. Veuillez répondre par 'oui' ou 'non'.")


def get_directories() -> list[str]:
    directories = []
    while True:
        path_to_discover = input("Donnez le chemin d'accès où chercher: ")
        if os.path.isdir(path_to_discover):
            directories.append(path_to_discover)
        else:
            print(f"Le chemin {path_to_discover} n'est pas valide.")

        end = input("Avez-vous terminé ? (oui/non): ").strip().lower()
        if end == "oui":
            break
    return directories

def file_contains_keyword_txt(path, keyword) -> bool:
    try:
        with open(path, "r", encoding='utf-8', errors='ignore') as file:
            for line in file:
                if keyword in line:
                    return True
    except Exception as e:
        print(f"Erreur {e} lors de la lecture du fichier {os.path.basename(path)}")
    return False

def file_contains_keyword_docx(path, keyword):
    doc = docx.Document(path)
    is_temp_file = (os.path.basename(path).startswith('~$') or '~$' in os.path.basename(path)
                    or '._' in os.path.basename(path))
    def read_all_doc_content(doc):
        fulltext = []
        for paragraph in doc.paragraphs:
            fulltext.append(paragraph.text)

        for tab in doc.tables:
            for rows in tab.rows:
                for cells in rows.cells:
                    fulltext.extend(read_all_doc_content(cells))
        return fulltext

    try:
        if is_temp_file or len(doc.paragraphs) <=0 :
            return False

        txt_content = read_all_doc_content(doc)

        for section in doc.sections:
            header = section.header
            footer = section.footer
            txt_content.extend(read_all_doc_content(header))
            txt_content.extend(read_all_doc_content(footer))

        for txt in txt_content:
            if keyword.lower() in txt.lower():
                return True
    except Exception as e:
        print(f"Erreur {e}")

    return False

def file_contain_keyword_doc(path, keyword):
    is_temp_file = (os.path.basename(path).startswith('~$') or '~$' in os.path.basename(path)
                    or '._' in os.path.basename(path))
    try:
        if is_temp_file:
            return False

        word = win32com.client.Dispatch('Word.Application')
        doc = word.Documents.Open(f'{path}')
        full_text = [doc.Content.Text]

        for table in doc.Tables:
            table_data = []
            for row in table.Rows:
                row_data = []
                for cell in row.Cells:
                    row_data.append(cell.Range.Text.strip())
                table_data.extend(row_data)
            full_text.append('\n'.join(['\t'.join(row) for row in table_data]))

        ##doc.Close()
       ##word.Quit()
        text_content = '\n\n'.join(full_text)

        if keyword.lower() in text_content.lower():
            return True

    except Exception as e:
        print(f"Erreur {e}")

    return False

def file_contain_keyword_xlsx(path, keyword):
    try:

        workbook = openpyxl.load_workbook(path)

        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows(values_only=True):
                for cell in row:
                    cell_text = str(cell).lower() if cell is not None else ''
                    if keyword.lower() in cell_text.lower():
                        return True
    except Exception as e:
        print(f"Erreur {e}")

    return False


def file_contain_keyword_xls(path, keyword):

    try:
        workbook = xlrd.open_workbook(path)
        for sheet in workbook.sheets():
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                for cell in row:
                    cell_text = str(cell.value).lower() if cell.value is not None else ''
                    if keyword.lower() in cell_text.lower():
                        return True

    except Exception as e:
        print(f"Erreur {e}")
    return False

def file_contain_keyword_pptx(path, keyword):
    keyword = keyword.lower()
    is_temp_file = (os.path.basename(path).startswith('~$') or '~$' in os.path.basename(path)
                    or '._' in os.path.basename(path))
    try:
        if is_temp_file:
            return False
        if not os.path.exists(path):
            print(f"Erreur: Le fichier n'existe pas à l'emplacement spécifié : {path}")
            return False

        presentation = Presentation(path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.lower()
                    if keyword in text:
                        return True

                if shape.shape_type == 19:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.lower()
                            if keyword in cell_text:
                                return True
    except Exception as e:
        print(f"Erreur {e}")
    return False

def file_contain_keyword_pdf(path, keyword):
    keyword = keyword.lower()
    is_temp_file = (os.path.basename(path).startswith('~$') or '~$' in os.path.basename(path)
                    or '._' in os.path.basename(path))
    try:
        if is_temp_file:
            return False
            # text content
        with open(path, 'rb') as file:
            reader = PdfReader(file)
            num_pages = len(reader.pages)
            for page_num in range(num_pages):
                page = reader.pages[page_num]
                text = page.extract_text()
                if text and keyword in text.lower():
                    return True

            # #  images
            # pages = convert_from_path(path)
            # for page in pages:
            #     text = pytesseract.image_to_string(page)
            #     if text and keyword in text.lower():
            #         return True

            # tables
            # tables = camelot.read_pdf(path, pages='all')
            # for table in tables:
            #     df = table.df  # Convert the table to a DataFrame
            #     for row in df.iterrows():
            #         if any(keyword in str(cell).lower() for cell in row[1]):
            #             return True
    except Exception as e:
        print(f"Erreur {e}")

    return False

def file_contain_keyword_ppt(path, keyword):
    keyword = keyword.lower()

    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(path, WithWindow=False)
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text = shape.TextFrame.TextRange.Text.lower()
                    if keyword in text:
                        ##presentation.Close()
                        ##powerpoint.Quit()
                        return True

                if shape.Type == 7:
                    table = shape.Table
                    for row in table.Rows:
                        for cell in row.Cells:
                            cell_text = cell.Shape.TextFrame.TextRange.Text.lower()
                            if keyword in cell_text:
                                ##presentation.Close()
                                ##powerpoint.Quit()
                                return True

        ##presentation.Close()
        ##powerpoint.Quit()
    except Exception as e:
        print(f"Erreur {e}")
    return False

def file_contain_keyword_visio(path, keyword):
    keyword = keyword.lower()

    try:
        visio = comtypes.client.CreateObject('Visio.Application')
        doc = visio.Documents.Open(path)
        for page in doc.Pages:
            for shape in page.Shapes:
                if shape.Text and keyword in shape.Text.lower():
                    doc.Close()
                    visio.Quit()
                    return True
        doc.Close()
        visio.Quit()
    except Exception as e:
        print(f"Erreur {e}")
    return False


def process_excel_doc(ext, path, keyword):
    if ext == 'xlsx':
        return file_contain_keyword_xlsx(path, keyword)
    elif ext == 'xls':
        return file_contain_keyword_xls(path, keyword)

def process_word_doc(ext , path, keyword):
    if ext == 'docx':
        return file_contains_keyword_docx(path, keyword)
    elif ext == 'doc':
        return file_contain_keyword_doc(path, keyword)
def process_powerpoint_doc(ext , path, keyword):
    if ext == 'pptx':
        return file_contain_keyword_pptx(path, keyword)
    elif ext == 'ppt':
        return file_contain_keyword_ppt(path, keyword)


