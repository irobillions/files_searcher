from abc import ABC, abstractmethod
import os
import textract
from pptx import Presentation
import PyPDF2
import docx
import pandas as pd
from transformers import pipeline
import xlrd
import win32com.client

# Initialize the summarizer pipeline
summarizer = pipeline("summarization")

class BaseSummarizer(ABC):
    @abstractmethod
    def summarize(self, file_path):
        pass

def summarize_text(text):
    max_chunk_size = 1024
    text_chunks = [text[i:i + max_chunk_size] for i in range(0, len(text), max_chunk_size)]
    summaries = []
    for chunk in text_chunks:
        summary = summarizer(chunk, max_length=150, min_length=30, do_sample=False)[0]['summary_text']
        summaries.append(summary)
    return " ".join(summaries)

class TextSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        summary = summarize_text(content)
        return f"TXT File Summary:\n{summary}"

class DocxSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        doc = docx.Document(file_path)
        content = '\n'.join([para.text for para in doc.paragraphs])
        summary = summarize_text(content)
        return f"DOCX File Summary:\n{summary}"

class DocSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(file_path)
        content = doc.Content.Text
        doc.Close()
        word.Quit()
        summary = summarize_text(content)
        return f"DOC File Summary:\n{summary}"

class XlsxSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        xls = pd.ExcelFile(file_path)
        sheet_descriptions = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            sheet_descriptions.append(f"Sheet: {sheet_name}\n{df.head().to_string(index=False)}\n")
        combined_sheets = "\n".join(sheet_descriptions)
        summary = summarize_text(combined_sheets)
        return f"XLSX File Summary:\n{summary}"

class XlsSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        xls = xlrd.open_workbook(file_path)
        sheet_descriptions = []
        for sheet_name in xls.sheet_names():
            sheet = xls.sheet_by_name(sheet_name)
            rows = [sheet.row_values(rownum) for rownum in range(min(sheet.nrows, 5))]
            df = pd.DataFrame(rows)
            sheet_descriptions.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n")
        combined_sheets = "\n".join(sheet_descriptions)
        summary = summarize_text(combined_sheets)
        return f"XLS File Summary:\n{summary}"

class PptxSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        prs = Presentation(file_path)
        slides_content = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_content.append("\n".join(slide_text))
        combined_slides = "\n".join(slides_content)
        summary = summarize_text(combined_slides)
        return f"PPTX File Summary:\n{summary}"

class PptSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        presentation = powerpoint.Presentations.Open(file_path, WithWindow=False)
        slides_content = []
        for slide in presentation.Slides:
            slide_text = []
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame") and shape.TextFrame.HasText:
                    slide_text.append(shape.TextFrame.TextRange.Text)
            slides_content.append("\n".join(slide_text))
        presentation.Close()
        powerpoint.Quit()
        combined_slides = "\n".join(slides_content)
        summary = summarize_text(combined_slides)
        return f"PPT File Summary:\n{summary}"

class PdfSummarizer(BaseSummarizer):
    def summarize(self, file_path):
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            content = []
            for page in reader.pages:
                content.append(page.extract_text())
        combined_content = "\n".join(content)
        summary = summarize_text(combined_content)
        return f"PDF File Summary:\n{summary}"
